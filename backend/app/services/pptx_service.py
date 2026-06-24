import json
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Brand colors
BLUE = RGBColor(37, 99, 235)
DARK = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(248, 250, 252)
MID_GRAY = RGBColor(100, 116, 139)


def _set_slide_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int, bold: bool = False,
                  color: RGBColor = DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def json_to_pptx(slides_json: str) -> bytes:
    """
    Convert JSON array of slide objects into a .pptx file.
    Returns pptx as bytes.
    """
    slides_data = json.loads(slides_json)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        is_cover = i == 0

        if is_cover:
            # Cover slide — dark blue background
            _set_slide_background(slide, BLUE)
            text_color = WHITE
            sub_color = RGBColor(186, 230, 253)
        else:
            # Content slides — white background
            _set_slide_background(slide, WHITE)
            text_color = DARK
            sub_color = MID_GRAY

        # Slide number pill (top left, skip on cover)
        if not is_cover:
            _add_text_box(
                slide,
                f"0{slide_data['slide_number']}" if slide_data['slide_number'] < 10 else str(slide_data['slide_number']),
                Inches(0.4), Inches(0.3), Inches(0.5), Inches(0.4),
                font_size=10, color=MID_GRAY
            )

        # Heading
        heading_top = Inches(1.8) if is_cover else Inches(1.0)
        _add_text_box(
            slide,
            slide_data.get("heading", slide_data.get("title", "")),
            Inches(0.8), heading_top, Inches(11.5), Inches(1.5),
            font_size=36 if is_cover else 28,
            bold=True,
            color=WHITE if is_cover else BLUE,
            align=PP_ALIGN.CENTER if is_cover else PP_ALIGN.LEFT
        )

        # Body content
        body_text = slide_data.get("body", "")
        if body_text:
            _add_text_box(
                slide,
                body_text,
                Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5),
                font_size=16,
                color=WHITE if is_cover else DARK,
                align=PP_ALIGN.CENTER if is_cover else PP_ALIGN.LEFT
            )

        # Speaker notes
        speaker_note = slide_data.get("speaker_note", "")
        if speaker_note:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_note

        # Bottom accent line on content slides
        if not is_cover:
            line = slide.shapes.add_connector(1, Inches(0.8), Inches(6.9), Inches(12.5), Inches(6.9))
            line.line.color.rgb = LIGHT_GRAY
            line.line.width = Pt(1)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()